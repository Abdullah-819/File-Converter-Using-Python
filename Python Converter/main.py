import os
import tkinter as tk
from tkinter import filedialog, ttk, messagebox
from tkinterdnd2 import DND_FILES, TkinterDnD
from PIL import Image, ImageTk, ImageDraw, ImageFont
from docx2pdf import convert as docx_to_pdf
from pdf2docx import Converter as pdf2docx_conv
from pdf2image import convert_from_path
from pptx import Presentation
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
import PyPDF2
import threading

# ===== Utility =====
def sanitize_text(text):
    return ''.join(ch for ch in text if ch.isprintable() or ch in '\n\r\t')

# ===== Conversion Functions =====
def convert_docx_to_pdf(input_path, output_path):
    docx_to_pdf(input_path, output_path)

def convert_pdf_to_docx(input_path, output_path, progress_callback=None):
    converter = pdf2docx_conv(input_path)
    num_pages = converter.doc.pages
    def update_progress(page):
        if progress_callback:
            percent = int((page/num_pages)*100)
            progress_callback(percent)
    converter.convert(output_path, start=0, end=num_pages, callback=update_progress)
    converter.close()

def convert_pdf_to_pptx(input_path, output_path, progress_callback=None):
    prs = Presentation()
    images = convert_from_path(input_path)
    total = len(images)
    for idx, img in enumerate(images):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        temp_img = f"temp_page_{idx}.jpg"
        img.save(temp_img)
        slide.shapes.add_picture(temp_img, 0, 0, width=prs.slide_width, height=prs.slide_height)
        os.remove(temp_img)
        if progress_callback:
            progress_callback(int((idx+1)/total*100))
    prs.save(output_path)

def convert_txt_to_pdf(input_path, output_path):
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(output_path)
    with open(input_path, "r", encoding="utf-8", errors="ignore") as f:
        text = sanitize_text(f.read())
    story = [Paragraph(text.replace("\n","<br/>"), styles["Normal"])]
    doc.build(story)

def convert_pdf_to_txt(input_path, output_path):
    reader = PyPDF2.PdfReader(open(input_path, "rb"))
    text = ""
    for page in reader.pages:
        page_text = page.extract_text() or ""
        text += sanitize_text(page_text)+"\n"
    with open(output_path, "w", encoding="utf-8", errors="ignore") as f:
        f.write(text)

def convert_docx_to_pptx(input_path, output_path):
    from docx import Document
    docx_file = Document(input_path)
    prs = Presentation()
    for para in docx_file.paragraphs:
        if para.text.strip():
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            if slide.placeholders:
                body = slide.placeholders[1]
                body.text_frame.text = sanitize_text(para.text)
            else:
                slide.shapes.add_textbox(0,0,prs.slide_width,prs.slide_height).text = sanitize_text(para.text)
    prs.save(output_path)

def convert_pptx_to_docx(input_path, output_path):
    from docx import Document
    prs = Presentation(input_path)
    docx_file = Document()
    for slide in prs.slides:
        if slide.shapes.title and slide.shapes.title.text.strip():
            docx_file.add_heading(sanitize_text(slide.shapes.title.text), level=1)
        for shape in slide.shapes:
            if hasattr(shape,"text") and shape.text.strip() and shape is not slide.shapes.title:
                docx_file.add_paragraph(sanitize_text(shape.text))
        docx_file.add_paragraph("\n")
    docx_file.save(output_path)

# ===== Circular Progress Bar =====
class CircularProgressBar(tk.Canvas):
    def __init__(self,parent,size=120,width=12,fg="#4CAF50",bg="#ddd",text_color="#333"):
        super().__init__(parent,width=size,height=size,bg=parent["bg"],highlightthickness=0)
        self.size=size
        self.width=width
        self.fg=fg
        self.bg=bg
        self.text_color=text_color
        self.create_oval(width//2,width//2,size-width//2,size-width//2,outline=self.bg,width=self.width,tags="bg_arc")
        self.arc = self.create_arc(width//2,width//2,size-width//2,size-width//2,start=-90,extent=0,
                                   style="arc",outline=self.fg,width=self.width,tags="progress_arc")
        self.text = self.create_text(size//2,size//2,text="0%",font=("Arial",18,"bold"),fill=self.text_color)

    def update_progress(self,percent):
        percent=max(0,min(100,percent))
        self.itemconfig(self.arc,extent=percent*3.6)
        self.itemconfig(self.text,text=f"{int(percent)}%")
        self.update_idletasks()

# ===== GUI =====
class FileConverterApp:
    def __init__(self, root):
        self.root=root
        self.root.title("Universal File Converter")
        self.root.geometry("900x550")
        self.root.resizable(False,False)
        self.file_path=None

        # Left frame
        self.left_frame=tk.Frame(root,width=450,height=550)
        self.left_frame.pack(side="left",fill="both")
        self.left_frame.pack_propagate(False)
        try:
            img=Image.open("background.png").resize((450,550),Image.LANCZOS)
        except:
            img=Image.new('RGB',(450,550),'#4A90E2')
            d=ImageDraw.Draw(img)
            font=ImageFont.load_default()
            d.text((50,250),"Universal\nConverter",fill=(255,255,255),font=font)
        self.bg_image=ImageTk.PhotoImage(img)
        self.bg_label=tk.Label(self.left_frame,image=self.bg_image)
        self.bg_label.pack(fill="both",expand=True)

        # Right frame
        self.right_frame=tk.Frame(root,width=450,height=550,bg="#f0f2f5",padx=20,pady=20)
        self.right_frame.pack(side="right",fill="both",expand=True)
        self.right_frame.pack_propagate(False)

        # Title
        self.title_label=tk.Label(self.right_frame,text="Universal File Converter",
                                  font=("Segoe UI",24,"bold"),bg="#f0f2f5",fg="#2C3E50")
        self.title_label.pack(pady=(20,15))

        # Drag & Drop
        self.drop_frame=tk.Frame(self.right_frame,bg="#e8edf2",highlightbackground="#cccccc",highlightthickness=2)
        self.drop_frame.pack(pady=10,fill="x",ipady=15)
        self.drop_label_text=tk.Label(self.drop_frame,text="Drag & Drop Your File Here\nor Click 'Upload File'",
                                      bg="#e8edf2",fg="#666",font=("Segoe UI",12,"italic"))
        self.drop_label_text.pack(expand=True)
        self.drop_frame.bind("<Button-1>", lambda e:self.upload_file())
        self.drop_label_text.bind("<Button-1>", lambda e:self.upload_file())
        root.drop_target_register(DND_FILES)
        root.dnd_bind('<<Drop>>', self.drop_file)
        self.drop_frame.dnd_bind('<<Drop>>',self.drop_file)
        self.drop_label_text.dnd_bind('<<Drop>>',self.drop_file)

        # File label
        self.file_label=tk.Label(self.right_frame,text="No file selected",font=("Segoe UI",10),bg="#f0f2f5",fg="#555")
        self.file_label.pack(pady=(5,15))

        # Conversion options
        self.conversion_var=tk.StringVar()
        self.dropdown=ttk.Combobox(self.right_frame,textvariable=self.conversion_var,state="readonly",width=30,font=("Segoe UI",11))
        self.dropdown['values']=[
            "DOCX → PDF",
            "PDF → DOCX",
            "PDF → PPTX",
            "TXT → PDF",
            "PDF → TXT",
            "DOCX → PPTX",
            "PPTX → DOCX"
        ]
        self.dropdown.pack(pady=10)
        self.dropdown.set("Select a conversion...")

        # Convert button
        self.convert_btn=ttk.Button(self.right_frame,text="🔄 Convert File",command=self.start_conversion_thread)
        self.convert_btn.pack(pady=20)

        # Circular progress
        self.circular_progress=CircularProgressBar(self.right_frame,size=120,bg="#f0f2f5")
        self.circular_progress.pack(pady=(0,10))

    def upload_file(self):
        path=filedialog.askopenfilename()
        if path:
            self.file_path=path
            self.file_label.config(text=f"Selected File: {os.path.basename(path)}")
            self.drop_label_text.config(text=f"File Loaded: {os.path.basename(path)}")

    def drop_file(self,event):
        self.file_path=event.data.strip("{}")
        if self.file_path:
            self.file_label.config(text=f"Selected File: {os.path.basename(self.file_path)}")
            self.drop_label_text.config(text=f"File Loaded: {os.path.basename(self.file_path)}")

    def start_conversion_thread(self):
        if not self.file_path:
            messagebox.showerror("Error","Please select a file first.")
            return
        if not self.conversion_var.get() or self.conversion_var.get()=="Select a conversion...":
            messagebox.showerror("Error","Please choose a conversion option.")
            return
        self.convert_btn.state(['disabled'])
        self.dropdown.config(state='disabled')
        threading.Thread(target=self._convert_file,daemon=True).start()

    def _convert_file(self):
        conv=self.conversion_var.get()
        ext_map={"DOCX → PDF":".pdf","PDF → DOCX":".docx","PDF → PPTX":".pptx",
                 "TXT → PDF":".pdf","PDF → TXT":".txt","DOCX → PPTX":".pptx","PPTX → DOCX":".docx"}
        initial_file_name=os.path.splitext(os.path.basename(self.file_path))[0]+ext_map.get(conv,"")
        save_path=filedialog.asksaveasfilename(defaultextension=ext_map.get(conv,""),
                                               initialfile=initial_file_name,parent=self.root)
        if not save_path:
            self._reset_gui()
            return
        try:
            # Call conversion functions with progress callback if supported
            if conv=="DOCX → PDF":
                convert_docx_to_pdf(self.file_path,save_path)
            elif conv=="PDF → DOCX":
                convert_pdf_to_docx(self.file_path,save_path,progress_callback=self.circular_progress.update_progress)
            elif conv=="PDF → PPTX":
                convert_pdf_to_pptx(self.file_path,save_path,progress_callback=self.circular_progress.update_progress)
            elif conv=="TXT → PDF":
                convert_txt_to_pdf(self.file_path,save_path)
            elif conv=="PDF → TXT":
                convert_pdf_to_txt(self.file_path,save_path)
            elif conv=="DOCX → PPTX":
                convert_docx_to_pptx(self.file_path,save_path)
            elif conv=="PPTX → DOCX":
                convert_pptx_to_docx(self.file_path,save_path)

            self.circular_progress.update_progress(100)
            messagebox.showinfo("Success",f"Conversion complete!\nSaved as:\n{save_path}")

        except Exception as e:
            messagebox.showerror("Error",str(e))
        finally:
            self._reset_gui()

    def _reset_gui(self):
        self.convert_btn.state(['!disabled'])
        self.dropdown.config(state='readonly')
        self.circular_progress.update_progress(0)

# ===== Run App =====
if __name__=="__main__":
    root=TkinterDnD.Tk()
    app=FileConverterApp(root)
    root.mainloop()
