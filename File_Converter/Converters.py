import os
from Utils import sanitize
from logger import log

from docx2pdf import convert as docx_to_pdf
from pdf2docx import Converter
from pdf2image import convert_from_path
from pptx import Presentation
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
import PyPDF2
# Converters.py (Corrected code)
from docx import Document

# ---------- DOCX → PDF ----------
def docx_to_pdf_conv(src, dst, callback):
    docx_to_pdf(src, dst)
    callback(100)
    log("DOCX → PDF completed")


# ---------- PDF → DOCX ----------
def pdf_to_docx_conv(src, dst, callback):
    # Use PyPDF2 to safely get the total page count
    try:
        reader = PyPDF2.PdfReader(src)
        total = len(reader.pages)
    except Exception as e:
        # Fallback to a single conversion step if PyPDF2 fails
        print(f"Warning: Could not get page count with PyPDF2. {e}")
        total = 1 
    
    cv = Converter(src)

    # Convert the first page to get initial setup done (this is a tricky hack)
    cv.convert(dst, start=0, end=1) 
    callback(1 / total * 100)
    
    # Now convert the rest of the pages (or just the whole document if total is 1)
    for page in range(1, total):
        cv.convert(dst, start=page, end=page+1)
        callback((page+1) / total * 100)

    cv.close()
    log("PDF → DOCX completed")

# ---------- PDF → PPTX ----------
def pdf_to_pptx_conv(src, dst, callback):
    prs = Presentation()
    pages = convert_from_path(src)

    total = len(pages)

    for i, img in enumerate(pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        temp = f"page_{i}.jpg"
        img.save(temp)

        slide.shapes.add_picture(temp, 0, 0, width=prs.slide_width)

        os.remove(temp)

        callback((i+1)/total * 100)

    prs.save(dst)
    log("PDF → PPTX completed")


# ---------- TXT → PDF ----------
def txt_to_pdf_conv(src, dst, callback):
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(dst)

    with open(src, "r", encoding="utf-8", errors="ignore") as f:
        text = sanitize(f.read())

    story = [Paragraph(text.replace("\n", "<br/>"), styles["Normal"])]
    doc.build(story)

    callback(100)
    log("TXT → PDF completed")


# ---------- PDF → TXT ----------
def pdf_to_txt_conv(src, dst, callback):
    reader = PyPDF2.PdfReader(src)
    total = len(reader.pages)
    text = ""

    for i, page in enumerate(reader.pages):
        text += (page.extract_text() or "") + "\n"
        callback((i+1)/total * 100)

    with open(dst, "w", encoding="utf-8") as f:
        f.write(sanitize(text))

    log("PDF → TXT completed")


# ---------- DOCX → PPTX ----------
def docx_to_pptx_conv(src, dst, callback):
    doc = Document(src)
    prs = Presentation()

    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    total = len(paras)

    for i, text in enumerate(paras):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = text[:100]
        callback((i+1)/total * 100)

    prs.save(dst)
    log("DOCX → PPTX completed")


# ---------- PPTX → DOCX ----------
def pptx_to_docx_conv(src, dst, callback):
    prs = Presentation(src)
    doc = Document()

    total = len(prs.slides)

    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                doc.add_paragraph(sanitize(shape.text))
        callback((i+1)/total * 100)

    doc.save(dst)
    log("PPTX → DOCX completed")
